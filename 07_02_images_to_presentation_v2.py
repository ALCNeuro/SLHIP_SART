#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Thu Mar 20 18:42:16 2025

@author: arthurlecoz

07_02_images_to_presentation.py
"""
# %% Paths & Packages

import os
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
import numpy as np
from glob import glob
from PIL import Image

# %% Functions 

# Helper function to get image dimensions in EMU (1 inch = 914400 EMU)
def image_dimensions_in_emu(image_path):
    im = Image.open(image_path)
    width_px, height_px = im.size
    # Get DPI if available; otherwise assume 96 dpi
    dpi = im.info.get('dpi', (96, 96))[0]
    width_in = width_px / dpi
    height_in = height_px / dpi
    emu_per_inch = 914400
    return int(width_in * emu_per_inch), int(height_in * emu_per_inch)

# Helper function to add a slide with an image at its native size (centered)
def add_slide_from_image(prs, fpath):
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        im = Image.open(fpath)
        width_px, height_px = im.size
        dpi = im.info.get('dpi', (96, 96))[0]
        width_in = width_px / dpi
        height_in = height_px / dpi
        emu_per_inch = 914400
        image_width = int(width_in * emu_per_inch)
        image_height = int(height_in * emu_per_inch)
        # Center the image on the slide
        left = int((prs.slide_width - image_width) / 2)
        top = int((prs.slide_height - image_height) / 2)
        slide.shapes.add_picture(fpath, left, top, width=image_width, height=image_height)
        
        # Extract subject code from file name (without extension)
        subject_code = os.path.splitext(os.path.basename(fpath))[0]
        comment_text = f"""{subject_code}

Commentaire libre:
...

...
Meilleure probabilité (W calme, W agité, N1, N2, N3, SP) :
...
Sur toute la fenêtre :

Sur les 5 seconds précédent la probe :

...
"""
        slide.notes_slide.notes_text_frame.text = comment_text
    except Exception as e:
        print(f"Error processing {fpath}: {e}")

# Helper function to add a blank slide with big centered text.
def add_text_slide(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    tf = txBox.text_frame
    tf.clear()
    # Create paragraph with centered, large text
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(48)
    p.alignment = 1  # Center alignment

# %% Paths and variables

if 'arthur' in os.getcwd():
    path_root = '/Volumes/DDE_ALC/PhD/SLHIP'
else:
    path_root = 'your_path'
path_scoring_coded = os.path.join(path_root, "07_Scoring", "Coded_figures")
this_identifiers = os.path.join(path_scoring_coded, 'identifiers.csv')

if os.path.exists(this_identifiers):
    df_id = pd.read_csv(this_identifiers)
    print("Identifiers loaded successfully.")

# Set the image folder and allowed extensions
img_folder = "/Volumes/DDE_ALC/PhD/SLHIP/07_Scoring/Coded_figures/what_is_this_090425"
allowed_extensions = (".jpg", ".jpeg", ".png", ".gif")

# Extract unique participant codes from filenames (assumes filename begins with the code)
codes_unique = np.unique([os.path.basename(f).split('_sess')[0]
                          for f in glob(os.path.join(img_folder, "*_session*"))])

# Define the range for the 40 probes (0 to 39)
probes_range = range(40)

# %%  Main Loop

for sub_id in codes_unique:
    for session in ["AM", "PM"]:
        # Gather all files matching participant and session
        these_files = glob(os.path.join(img_folder, f"{sub_id}*{session}*.png"))
        
        # Partition the files based on filename indicators:
        # RS_1: resting state before task (e.g., ...RS_1_w_0.png to ...RS_1_w_9.png)
        # RS_2: resting state after task (e.g., ...RS_2_w_0.png to ...RS_2_w_9.png)
        # Probes: files with "before_probe" or "during_probe"
        rs1_files = [f for f in these_files if "RS_1" in f]
        rs2_files = [f for f in these_files if "RS_2" in f]
        probe_before_files = [f for f in these_files if "before_probe" in f]
        probe_during_files = [f for f in these_files if "during_probe" in f]

        # Sort RS files by window index using a helper extraction function.
        def extract_rs_num(filename, marker):
            try:
                base = os.path.basename(filename)
                return int(base.split(f"{marker}_w_")[-1].split('.')[0])
            except Exception:
                return -1
        
        rs1_files = sorted(rs1_files, key=lambda f: extract_rs_num(f, "RS_1"))
        rs2_files = sorted(rs2_files, key=lambda f: extract_rs_num(f, "RS_2"))

        # For probes, order the slides as: before_probe_X then during_probe_X for each probe.
        probe_files = []
        for probe in probes_range:
            before_list = [f for f in probe_before_files if f"before_probe_{probe}" in f]
            if before_list:
                probe_files.append(before_list[0])
            during_list = [f for f in probe_during_files if f"during_probe_{probe}" in f]
            if during_list:
                probe_files.append(during_list[0])
        
        # Create a new presentation for this session.
        prs = Presentation()
        
        # Set the slide dimensions to that of the first image in this session (if available)
        if these_files:
            first_img_path = these_files[0]
            slide_width, slide_height = image_dimensions_in_emu(first_img_path)
            prs.slide_width = slide_width
            prs.slide_height = slide_height
        else:
            # Default size if no images are found
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        
        # ---------------------------------------------------------------------
        # Build the presentation in the desired order:
        # 1. A text slide marking the start of RS_1.
        # 2. RS_1 image slides.
        # 3. A text slide marking the transition between RS_1 and probes.
        # 4. Probe image slides.
        # 5. A text slide marking the start of RS_2.
        # 6. RS_2 image slides.
        # ---------------------------------------------------------------------
        add_text_slide(prs, "DEBUT DU RESTING STATE 1 :")
        
        for f in rs1_files:
            add_slide_from_image(prs, f)
            
        add_text_slide(prs, "FIN DU RESTING STATE 1, DÉBUT DES PROBES")
        
        for f in probe_files:
            add_slide_from_image(prs, f)
            
        add_text_slide(prs, "FIN DES PROBES, DÉBUT DU RESTING STATE 2")
        
        for f in rs2_files:
            add_slide_from_image(prs, f)
        
        # Save the presentation with a filename that indicates the participant and session.
        output_file = os.path.join(img_folder, f"{sub_id}_{session}_presentation.pptx")
        prs.save(output_file)
        print(f"Saved presentation for participant {sub_id} session {session}: {output_file}")
    
    
    
    
