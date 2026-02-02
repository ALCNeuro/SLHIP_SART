#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Tue Dec 16 11:55:15 2025

@author: arthurlecoz

07_03_presentation_to_scoring.py
"""
# %% Paths & Packages
import os
import re
from glob import glob
import numpy as np
import pandas as pd
from pptx import Presentation
from scipy.io import loadmat

if "arthur" in os.getcwd():
    path_root = "/Volumes/DDE_ALC/PhD/SLHIP"
else:
    path_root = "your_path"

path_scored = os.path.join(path_root, "07_Scoring", "Scored", "Isabelle")
path_scoring_coded = os.path.join(path_root, "07_Scoring", "Coded_figures")
path_figures = os.path.join(path_root, "07_Scoring", "Figures")
path_identifiers = os.path.join(path_scoring_coded, "identifiers.csv")
path_raw_experiment = os.path.join(path_root, "00_Raw", "experiment")

out_dir = path_scored
out_raw_csv = os.path.join(out_dir, "scoring_raw_audit.csv")
out_struct_csv = os.path.join(out_dir, "scoring_structured.csv")

# %% Constants & Dic

probe_col = [
    "nprobe","t_probe_th","t_probe_act","nblock","block_cond","ntrial",
    "PQ1_respkey","PQ2_respkey","PQ3_respkey",
    "PQ1_resptime","PQ2_resptime","PQ3_resptime",
    "PQ1_questime","PQ2_questime","PQ3_questime",
    "PQ1_respval","PQ2_respval","PQ3_respval"
]

ms_dic = {
    0 : "MISS",
    1 : "ON",
    2 : "MW",
    3 : "DISTRACTED",
    4 : "HALLU",
    5 : "MB",
    6 : "FORGOT"
}

# headers (case-insensitive, robust to accents/spacing)
H_COMMENT = re.compile(r"^\s*Commentaire\s+libre\s*:\s*$", re.IGNORECASE)
H_BESTPROB = re.compile(r"^\s*Meilleure\s+probabilit", re.IGNORECASE)  # match the long line
H_30S = re.compile(r"^\s*Sur\s+toute\s+la\s+fen", re.IGNORECASE)         # "fenêtre" sometimes without accent
H_5S = re.compile(r"^\s*Sur\s+les\s+5\s*second", re.IGNORECASE)

# subject_code pattern: 6 chars + _session_ + AM/PM + ....
SUBJECT_CODE_RE = re.compile(
    r"^(?P<rid>[A-Za-z0-9]{6})_session_(?P<session>AM|PM)_(?P<rest>.+)$"
)

# label tokens you want to keep "brut" (we'll store matches as they appear + a canonical version)
LABEL_RE = re.compile(
    r"\b(W\s*calme|W\s*agit[eé]|W|N1|N2|N3|SP)\b",
    re.IGNORECASE
)

def canonicalize_label(lbl: str) -> str:
    s = lbl.strip().lower()
    s = s.replace("é", "e").replace("è", "e").replace("ê", "e")
    s = re.sub(r"\s+", " ", s)
    if s == "w calme":
        return "W calme"
    if s in ["w agite", "w agité", "w agité"]:
        return "W agité"
    if s == "w":
        return "W"
    if s in ["n1", "n2", "n3"]:
        return s.upper()
    if s == "sp":
        return "SP"
    return lbl.strip()


# %% Fun Helpers

def first_nonempty_line(lines):
    for ln in lines:
        if ln.strip():
            return ln.strip()
    return None

def clean_block(lines):
    """
    Remove template artifacts while keeping "raw" clinician text.
    - drop empty lines
    - drop lines exactly "..."
    """
    keep = []
    for ln in lines:
        t = ln.strip()
        if not t:
            continue
        if t == "...":
            continue
        keep.append(ln.rstrip())
    return "\n".join(keep).strip() if keep else np.nan

def find_header_indices(lines):
    """
    Return dict of header->index (first occurrence).
    """
    idx = {"comment": None, "bestprob": None, "s30": None, "s5": None}
    for i, ln in enumerate(lines):
        if idx["comment"] is None and H_COMMENT.match(ln):
            idx["comment"] = i
        elif idx["bestprob"] is None and H_BESTPROB.match(ln):
            idx["bestprob"] = i
        elif idx["s30"] is None and H_30S.match(ln):
            idx["s30"] = i
        elif idx["s5"] is None and H_5S.match(ln):
            idx["s5"] = i
    return idx

def extract_block_between(lines, start_i, end_i):
    if start_i is None:
        return []
    start = start_i + 1
    end = end_i if end_i is not None else len(lines)
    if start >= end:
        return []
    return lines[start:end]

def extract_tokens(text):
    if not isinstance(text, str) or not text.strip():
        return [], []
    raw = [m.group(0) for m in LABEL_RE.finditer(text)]
    canon = [canonicalize_label(x) for x in raw]
    return raw, canon

def parse_subject_code(subject_code):
    """
    Returns dict with id6char, session_time, recording_type, idx (probe or window), plus flags.
    """
    out = {
        "id6char": np.nan,
        "session_time": np.nan,
        "recording_type": np.nan,
        "probe_idx": np.nan,
        "w_idx": np.nan,
        "subject_code_ok": False
    }
    if not subject_code:
        return out

    m = SUBJECT_CODE_RE.match(subject_code.strip())
    if not m:
        return out

    rid = m.group("rid")
    session = m.group("session")
    rest = m.group("rest")

    out["id6char"] = rid
    out["session_time"] = session
    out["subject_code_ok"] = True

    if "RS_1" in rest:
        out["recording_type"] = "RS1"
        wm = re.search(r"_w_(\d+)\b", rest)
        if wm:
            out["w_idx"] = int(wm.group(1))
    elif "RS_2" in rest:
        out["recording_type"] = "RS2"
        wm = re.search(r"_w_(\d+)\b", rest)
        if wm:
            out["w_idx"] = int(wm.group(1))
    elif "before_probe" in rest:
        out["recording_type"] = "Probe_before"
        pm = re.search(r"_probe_(\d+)\b", rest)
        if pm:
            out["probe_idx"] = int(pm.group(1))
    elif "during_probe" in rest:
        out["recording_type"] = "Probe_during"
        pm = re.search(r"_probe_(\d+)\b", rest)
        if pm:
            out["probe_idx"] = int(pm.group(1))
    else:
        out["recording_type"] = "UNKNOWN"

    return out

def read_notes_text(slide):
    """
    Return full notes text or '' if missing.
    """
    try:
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        if tf is None:
            return ""
        return tf.text or ""
    except Exception:
        return ""

def get_mindstate_for_probe(sub_id, session_time, probe_idx, mindstate_cache):
    """
    Load .mat once per (sub_id, session_time) and return mindstate for probe_idx.
    """
    if pd.isna(probe_idx):
        return np.nan

    key = (sub_id, session_time)
    if key not in mindstate_cache:
        # Load the appropriate .mat
        behav_paths = glob(os.path.join(path_raw_experiment, f"sub_{sub_id}", "*.mat"))
        if len(behav_paths) < 1:
            mindstate_cache[key] = None
        else:
            # Prefer file containing AM/PM in filename if possible
            preferred = [p for p in behav_paths if f"_{session_time}" in os.path.basename(p)]
            if preferred:
                behav_path = preferred[0]
            else:
                # fallback to old assumption: AM=first, PM=second if 2 files
                behav_paths = sorted(behav_paths)
                if session_time == "AM":
                    behav_path = behav_paths[0]
                else:
                    behav_path = behav_paths[1] if len(behav_paths) > 1 else behav_paths[0]

            try:
                mat = loadmat(behav_path)
                df_probe = pd.DataFrame(mat["probe_res"], columns=probe_col)
                if any(df_probe.PQ1_respval.isna()):
                    df_probe.PQ1_respval.replace(np.nan, 0, inplace=True)
                ms_answers = np.array([ms_dic.get(int(v), "UNK") for v in df_probe.PQ1_respval.values])
                mindstate_cache[key] = ms_answers
            except Exception:
                mindstate_cache[key] = None

    ms_answers = mindstate_cache.get(key, None)
    if ms_answers is None:
        return np.nan
    try:
        probe_idx = int(probe_idx)
        if probe_idx < 0 or probe_idx >= len(ms_answers):
            return np.nan
        return ms_answers[probe_idx]
    except Exception:
        return np.nan


# %% Load Identifiers

if not os.path.exists(path_identifiers):
    raise FileNotFoundError(f"Missing identifiers.csv at: {path_identifiers}")

df_id = pd.read_csv(path_identifiers)
# Expect columns: sub_id, random_id
if not {"sub_id", "random_id"}.issubset(df_id.columns):
    raise ValueError("identifiers.csv must contain columns: sub_id, random_id")

rid_to_sub = dict(
    zip(df_id["random_id"].astype(str), df_id["sub_id"].astype(str))
    )


# %% Main script

pptx_files = sorted(glob(os.path.join(path_scored, "*.pptx")))
if len(pptx_files) == 0:
    raise FileNotFoundError(f"No .pptx found in: {path_scored}")

rows_raw = []
rows_struct = []
mindstate_cache = {}

for pptx_path in pptx_files:
    prs = Presentation(pptx_path)

    for slide_idx, slide in enumerate(prs.slides):
        notes = read_notes_text(slide)
        lines = notes.replace("\r\n", "\n").replace("\r", "\n").split("\n")

        subject_code = first_nonempty_line(lines)
        meta = parse_subject_code(subject_code)

        # If slide is a marker ("DEBUT..." etc.), subject_code won't match -> keep raw row but flag and skip structured
        flags = []
        if not meta["subject_code_ok"]:
            if subject_code and any(x in subject_code.upper() for x in ["DEBUT", "FIN"]):
                flags.append("marker_slide")
            else:
                flags.append("missing_or_invalid_subject_code")

        # Parse headers on remaining lines (after first non-empty line)
        # We keep everything, but try to interpret blocks.
        idxs = find_header_indices(lines)

        # Determine block boundaries by taking the next header index after each header
        header_positions = sorted([i for i in idxs.values() if i is not None])
        def next_header(after_i):
            for hp in header_positions:
                if hp > after_i:
                    return hp
            return None

        comment_lines = extract_block_between(lines, idxs["comment"], next_header(idxs["comment"]) if idxs["comment"] is not None else None)
        bestprob_lines = extract_block_between(lines, idxs["bestprob"], next_header(idxs["bestprob"]) if idxs["bestprob"] is not None else None)
        s30_lines = extract_block_between(lines, idxs["s30"], next_header(idxs["s30"]) if idxs["s30"] is not None else None)
        s5_lines = extract_block_between(lines, idxs["s5"], next_header(idxs["s5"]) if idxs["s5"] is not None else None)

        # Clean blocks (drop "..." and empty) but keep raw clinician text
        commentaire = clean_block(comment_lines)
        bestprob = clean_block(bestprob_lines)
        scoring30s = clean_block(s30_lines)
        scoring5s = clean_block(s5_lines)

        # Tokens extraction (raw + canonical)
        all_raw_tokens, all_canon_tokens = extract_tokens(notes)
        b_raw_tokens, b_canon_tokens = extract_tokens(bestprob if isinstance(bestprob, str) else "")
        s30_raw_tokens, s30_canon_tokens = extract_tokens(scoring30s if isinstance(scoring30s, str) else "")
        s5_raw_tokens, s5_canon_tokens = extract_tokens(scoring5s if isinstance(scoring5s, str) else "")

        # Flags about missing headers / suspicious placement
        if idxs["comment"] is None:
            flags.append("missing_header_commentaire")
        if idxs["s30"] is None:
            flags.append("missing_header_30s")
        if idxs["s5"] is None:
            flags.append("missing_header_5s")
        if idxs["bestprob"] is None:
            flags.append("missing_header_bestprob")

        if len(s30_canon_tokens) > 1:
            flags.append("multiple_labels_30s")
        if len(s5_canon_tokens) > 1:
            flags.append("multiple_labels_5s")
        if len(all_canon_tokens) == 0 and meta["subject_code_ok"]:
            flags.append("no_labels_anywhere")

        # If no scoring in blocks but labels exist elsewhere, likely answered outside expected place
        if (not isinstance(scoring30s, str) or not scoring30s.strip()) and len(all_canon_tokens) > 0:
            flags.append("labels_outside_30s_block")
        if (not isinstance(scoring5s, str) or not scoring5s.strip()) and len(all_canon_tokens) > 0:
            flags.append("labels_outside_5s_block")

        # Map random_id -> sub_id/group
        rid = meta["id6char"]
        sub_id = rid_to_sub.get(str(rid), np.nan) if meta["subject_code_ok"] else np.nan
        group = sub_id.split("_")[0] if isinstance(sub_id, str) and "_" in sub_id else np.nan

        # Mindstate for probes only
        mindstate = np.nan
        if meta["recording_type"] in ["Probe_before", "Probe_during"] and isinstance(sub_id, str) and not pd.isna(meta["probe_idx"]):
            mindstate = get_mindstate_for_probe(sub_id, meta["session_time"], meta["probe_idx"], mindstate_cache)

        # RAW AUDIT ROW (keep even for marker slides)
        rows_raw.append({
            "pptx_file": os.path.basename(pptx_path),
            "pptx_path": pptx_path,
            "slide_index": slide_idx,
            "subject_code_first_line": subject_code,
            "subject_code_ok": meta["subject_code_ok"],
            "id6char": meta["id6char"],
            "session_time": meta["session_time"],
            "recording_type": meta["recording_type"],
            "probe_idx": meta["probe_idx"],
            "w_idx": meta["w_idx"],
            "sub_id": sub_id,
            "group": group,
            "mindstate": mindstate,
            "raw_notes_text": notes,
            "commentaire_block_clean": commentaire,
            "bestprob_block_clean": bestprob,
            "scoring30s_block_clean": scoring30s,
            "scoring5s_block_clean": scoring5s,
            "all_label_tokens_raw": "|".join(all_raw_tokens) if all_raw_tokens else "",
            "all_label_tokens_canon": "|".join(all_canon_tokens) if all_canon_tokens else "",
            "bestprob_tokens_raw": "|".join(b_raw_tokens) if b_raw_tokens else "",
            "bestprob_tokens_canon": "|".join(b_canon_tokens) if b_canon_tokens else "",
            "scoring30s_tokens_raw": "|".join(s30_raw_tokens) if s30_raw_tokens else "",
            "scoring30s_tokens_canon": "|".join(s30_canon_tokens) if s30_canon_tokens else "",
            "scoring5s_tokens_raw": "|".join(s5_raw_tokens) if s5_raw_tokens else "",
            "scoring5s_tokens_canon": "|".join(s5_canon_tokens) if s5_canon_tokens else "",
            "flags": "|".join(flags) if flags else ""
        })

        # STRUCTURED ROW: only if it looks like a real scoring slide (valid subject_code)
        if meta["subject_code_ok"]:
            rows_struct.append({
                "id6char": meta["id6char"],
                "sub_id": sub_id,
                "group": group,
                "recording_type": meta["recording_type"],
                "session_time": meta["session_time"],
                "probe_idx": meta["probe_idx"],
                "w_idx": meta["w_idx"],
                "commentaire": commentaire,
                "bestprob": bestprob,
                "scoring30s": scoring30s,
                "scoring5s": scoring5s,
                "mindstate": mindstate,
                "pptx_file": os.path.basename(pptx_path),
                "slide_index": slide_idx,
                # Tokens (useful for manual validation / later normalization)
                "scoring30s_tokens_raw": "|".join(s30_raw_tokens) if s30_raw_tokens else "",
                "scoring5s_tokens_raw": "|".join(s5_raw_tokens) if s5_raw_tokens else "",
                "all_label_tokens_raw": "|".join(all_raw_tokens) if all_raw_tokens else "",
                "flags": "|".join(flags) if flags else ""
            })


# Export Csv
df_raw = pd.DataFrame(rows_raw)
df_struct = pd.DataFrame(rows_struct)

# Sort nicely (structured)
sort_cols = ["sub_id", "session_time", "recording_type", "probe_idx", "w_idx", "pptx_file", "slide_index"]
sort_cols = [c for c in sort_cols if c in df_struct.columns]
df_struct = df_struct.sort_values(sort_cols).reset_index(drop=True)

df_raw.to_csv(out_raw_csv, index=False)
df_struct.to_csv(out_struct_csv, index=False)

print(f"✅ Raw audit CSV saved:       {out_raw_csv}")
print(f"✅ Structured scoring CSV saved: {out_struct_csv}")
print(f"Rows (raw): {len(df_raw)} | Rows (structured): {len(df_struct)}")

# %% Adjust Scoring

finegrained_30_dic = {
    'w'         : 'W', 
    'N1'        : 'N1', 
    'W'         : 'W', 
    'W|N1'      : 'W_N1', 
    'SP'        : 'R', 
    'W calme'   : 'W', 
    'N1|W'      : 'W_N1', 
    'W|SP'      : 'W_R',
    'N1|N1|W'   : 'W_N1',  
    'N2'        : 'N2', 
    'W|SP|W'    : 'W_R', 
    'N1|W|N1'   : 'W_N1',  
    'SP|W'      : 'W_R', 
    'W|N1|W'    : 'W_N1', 
    'n1|W'      : 'W_N1',
    'N1|N1'     : 'N1', 
    'W|N1|W|N1' : 'W_N1',
    ""          : np.nan
    }

classic_30_dic = {
    'w'         : 'W', 
    'N1'        : 'N1', 
    'W'         : 'W', 
    'W|N1'      : 'N1', 
    'SP'        : 'R', 
    'W calme'   : 'W', 
    'N1|W'      : 'N1', 
    'W|SP'      : 'R',
    'N1|N1|W'   : 'N1',  
    'N2'        : 'N2', 
    'W|SP|W'    : 'R', 
    'N1|W|N1'   : 'N1',  
    'SP|W'      : 'R', 
    'W|N1|W'    : 'N1', 
    'n1|W'      : 'N1',
    'N1|N1'     : 'N1', 
    'W|N1|W|N1' : 'N1',
    ""          : np.nan,
    np.nan      : np.nan
    }

classic_5_dic = {
    ''          : '', 
    'w'         : 'W', 
    'W calme'   : 'W', 
    'N1'        : 'N1', 
    'W'         : 'W', 
    'N1|N2|N2'  : 'N2', 
    'N1|W'      : 'W', 
    'SP'        : 'R', 
    'N2'        : 'N2'
    }

"""
Si je n'ai pas de scoring sur les 5 dernières secondes, 
inférer le dernier indiqué par Isabelle
Si vide et 30s : "W|N1" -> choisir N1
"""

# %% 

simple_scoring = []
fined_scoring = []

for i, score in enumerate(df_struct.scoring30s_tokens_raw) :
    simple_scoring.append(classic_30_dic[score])
    fined_scoring.append(finegrained_30_dic[score])
    
# %% 

df_struct.insert(
    loc = 13,
    column = "tradi_scoring", 
    value = simple_scoring)

df_struct.insert(
    loc = 13,
    column = "finer_scoring", 
    value = simple_scoring)

out_scoring_adapted = os.path.join(out_dir, "scoring_structured_classic_stages.csv")
df_struct.to_csv(out_scoring_adapted)

# %% 30s per stage per subject per recording type per session

cols = ["sub_id", "group", "recording_type", "session_time", "tradi_scoring", "percentage"]

dic_oi = {col : [] for col in cols}

for i, sub_id in enumerate(df_struct.sub_id.unique()) :
    sub_df = df_struct.loc[df_struct.sub_id == sub_id]
    group = sub_id[:2]
    
    for recording_type in sub_df.recording_type.unique() :
        for session_time in sub_df.session_time.unique() :
            
            this_df = sub_df.loc[
                (sub_df.recording_type == recording_type) 
                & (sub_df.session_time == session_time)
                ]
            
            if not len(this_df) : continue
            
            for stage in ["W", "R", "N1", "N2"] :
                
                percentage = len(this_df.loc[this_df.tradi_scoring == stage])/len(this_df) * 100
                
                dic_oi["sub_id"].append(sub_id)
                dic_oi["group"].append(group)
                dic_oi["recording_type"].append(recording_type)
                dic_oi["session_time"].append(session_time)
                dic_oi["tradi_scoring"].append(stage)
                dic_oi["percentage"].append(percentage)

df_30s_classic = pd.DataFrame.from_dict(dic_oi)    
this_savepath = os.path.join(path_figures, "table_percentage_recording_time_tradi.csv")
df_30s_classic.to_csv(this_savepath)                 

# %% Probe locked percentages with mindstate                
                
cols = ["sub_id", "group", "mindstate", "session_time", "tradi_scoring", "percentage"]
mindstates = ["ON", "MW", "MB", "HALLU", "FORGOT", "MISS", "DISTRACTED"]

dic_oi = {col : [] for col in cols}

for i, sub_id in enumerate(df_struct.sub_id.unique()) :
    sub_df = df_struct.loc[
        (df_struct.sub_id == sub_id)
        & (df_struct.recording_type == "Probe_during")
        ]
    
    group = sub_id[:2]
    
    for mindstate in mindstates :
        for session_time in sub_df.session_time.unique() :
            
            this_df = sub_df.loc[
                (sub_df.mindstate == mindstate) 
                & (sub_df.session_time == session_time)
                ]
            
            if not len(this_df) : continue
            
            for stage in ["W", "R", "N1", "N2"] :
                
                percentage = len(this_df.loc[this_df.tradi_scoring == stage])/len(this_df) * 100
                
                dic_oi["sub_id"].append(sub_id)
                dic_oi["group"].append(group)
                dic_oi["mindstate"].append(mindstate)
                dic_oi["session_time"].append(session_time)
                dic_oi["tradi_scoring"].append(stage)
                dic_oi["percentage"].append(percentage)

df_mindstate_classic = pd.DataFrame.from_dict(dic_oi)    
this_savepath = os.path.join(path_figures, "table_percentage_mindstate_time_tradi.csv")
df_mindstate_classic.to_csv(this_savepath)                
    
    
