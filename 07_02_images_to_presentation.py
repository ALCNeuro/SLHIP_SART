#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Thu Mar 20 18:42:16 2025

@author: arthurlecoz

07_02_images_to_presentation.py
"""
# %% 
import os
import random
from pptx import Presentation
from pptx.util import Inches

# Specify the folder containing your images
img_folder = "/Volumes/DDE_ALC/PhD/SLHIP/07_Scoring/Coded_figures/what_is_this_080425"
# Define allowed image extensions
allowed_extensions = (".jpg", ".jpeg", ".png", ".gif")

# Retrieve and shuffle image filenames
images = [f for f in os.listdir(img_folder) if f.lower().endswith(allowed_extensions)]
random.shuffle(images)

# Set batch size for each presentation
batch_size = 100

# Process images in batches
total_images = len(images)
for batch_start in range(0, total_images, batch_size):
    prs = Presentation()
    batch_images = images[batch_start:batch_start + batch_size]
    
    for image_file in batch_images:
        # Add a blank slide (layout index 6 typically corresponds to a blank slide)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img_path = os.path.join(img_folder, image_file)
        # Insert the image so that it fills the entire slide
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), 
                                   width=prs.slide_width, height=prs.slide_height)
        
        # Extract subject code by removing the file extension
        subject_code = os.path.splitext(image_file)[0]
        # Create the comment text with the subject code included
        comment_text = f"""{subject_code}

Commentaire libre:
...

...    
Meilleure probabilité (W calme, W agité, N1, N2, N3, SP) :
...

... 
"""
        # Add the comment in the slide notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = comment_text
    
    # Save the current batch presentation with a unique filename
    batch_number = batch_start // batch_size + 1
    prs.save(f"/Volumes/DDE_ALC/PhD/SLHIP/07_Scoring/Presentations/what_is_it_{batch_number}.pptx")
    print(f"Saved my_presentation_batch_{batch_number}.pptx with {len(batch_images)} images.")
    
# %% clean up
import pandas as pd
import shutil
df = pd.read_csv('/Volumes/DDE_ALC/PhD/SLHIP/07_Scoring/Coded_figures/code_correspondance_N1.csv')
destination_dir = '/Volumes/DDE_ALC/PhD/SLHIP/07_Scoring/Coded_figures/to_erase'

temp = []
for image in images :
    if image[:-4] not in df.code.values : 
        temp.append(os.path.join(img_folder, image))

for filename in temp :
    destination_file = os.path.join(destination_dir, filename.split('figures/')[-1])
    try:
        shutil.move(filename, destination_file)
        print(f"Moved file {filename} to {destination_dir}")
    except Exception as e:
        print(f"Error moving file {filename}: {e}")
